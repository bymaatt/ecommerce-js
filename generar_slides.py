from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ===== Paleta de colores del proyecto =====
NAVY      = RGBColor(0x1a, 0x1a, 0x2e)   # fondo oscuro / header
BLUE      = RGBColor(0x0f, 0x34, 0x60)   # botones / primario
RED       = RGBColor(0xe9, 0x45, 0x60)   # acento rojo
WHITE     = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG  = RGBColor(0xf5, 0xf5, 0xf5)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE= RGBColor(0x1a, 0x4a, 0x8a)
CODE_BG   = RGBColor(0x0d, 0x1b, 0x2a)   # fondo de codigo oscuro

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout completamente en blanco

# =============================================
#  Helpers
# =============================================

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox

def add_multiline_text(slide, lines, x, y, w, h,
                       font_size=16, bold=False, color=WHITE,
                       align=PP_ALIGN.LEFT, line_spacing=1.15):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return txBox

def slide_dark_header(slide, title, subtitle=None):
    """Banda superior oscura con título y subtítulo opcional."""
    add_rect(slide, 0, 0, 13.33, 1.5, NAVY)
    add_rect(slide, 0, 1.5, 13.33, 0.08, RED)
    add_text_box(slide, title, 0.5, 0.2, 12, 0.9,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.5, 0.95, 12, 0.55,
                     font_size=15, bold=False, color=RGBColor(0xcc, 0xcc, 0xdd),
                     align=PP_ALIGN.LEFT)

def add_bullet_box(slide, bullets, x, y, w, h,
                   font_size=15, color=RGBColor(0x22, 0x22, 0x44), dot_color=None):
    """Caja con viñetas •  texto."""
    if dot_color is None:
        dot_color = RED
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # punto de color
        r1 = p.add_run()
        r1.text = "● "
        r1.font.size = Pt(font_size - 1)
        r1.font.color.rgb = dot_color
        r1.font.name = "Segoe UI"
        # texto
        r2 = p.add_run()
        r2.text = bullet
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color
        r2.font.name = "Segoe UI"
    return txBox

def add_code_block(slide, code_lines, x, y, w, h, font_size=10):
    """Bloque de código estilo dark."""
    bg = add_rect(slide, x, y, w, h, CODE_BG)
    # borde izquierdo de acento
    add_rect(slide, x, y, 0.07, h, RED)
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.15),
        Inches(w - 0.25), Inches(h - 0.3)
    )
    txBox.word_wrap = False
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0xc8, 0xe0, 0xff)
        run.font.name = "Consolas"
    return txBox

def add_tag_label(slide, text, x, y, color=BLUE):
    """Etiqueta de chip/tag redondeada (simulada con rectángulo)."""
    add_rect(slide, x, y, len(text) * 0.12 + 0.2, 0.32, color)
    add_text_box(slide, text, x + 0.07, y + 0.03, len(text) * 0.12 + 0.1, 0.28,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def add_card(slide, title, lines, x, y, w, h,
             title_color=NAVY, bg=WHITE, accent=RED, font_size=13):
    """Tarjeta con título y contenido."""
    add_rect(slide, x, y, w, h, bg)
    add_rect(slide, x, y, w, 0.06, accent)
    add_text_box(slide, title, x + 0.15, y + 0.1, w - 0.3, 0.38,
                 font_size=14, bold=True, color=title_color)
    for i, line in enumerate(lines):
        add_text_box(slide, line, x + 0.15, y + 0.48 + i * 0.32, w - 0.3, 0.35,
                     font_size=font_size, color=GRAY_TEXT)

# =============================================
#  SLIDE 1 — PORTADA
# =============================================
slide1 = prs.slides.add_slide(BLANK)

# Fondo completo oscuro
add_rect(slide1, 0, 0, 13.33, 7.5, NAVY)

# Decoración: banda roja vertical lateral
add_rect(slide1, 0, 0, 0.18, 7.5, RED)

# Decoración: bloque azul
add_rect(slide1, 0.18, 0, 13.15, 0.1, BLUE)

# Logo / Icono (emoji grande simulado con texto)
add_text_box(slide1, "🛒", 1.5, 1.1, 2.5, 2.0,
             font_size=72, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# Título principal
add_text_box(slide1, "TechStore", 4.2, 1.2, 8.5, 1.4,
             font_size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Línea roja decorativa
add_rect(slide1, 4.2, 2.7, 5.5, 0.1, RED)

# Subtítulo
add_text_box(slide1, "Simulador de Compras", 4.2, 2.85, 8.5, 0.65,
             font_size=26, bold=False, color=RGBColor(0xcc, 0xcc, 0xdd),
             align=PP_ALIGN.LEFT)

# Chips de tecnologías
for i, (tag, color) in enumerate([
        ("HTML5", RGBColor(0xe3, 0x4f, 0x26)),
        ("CSS3",  RGBColor(0x15, 0x72, 0xb6)),
        ("JavaScript", RGBColor(0xf0, 0xdb, 0x4f)),
        ("Toastify JS", RED),
        ("Fetch API", BLUE),
]):
    add_rect(slide1, 4.2 + i * 1.82, 3.65, 1.65, 0.38, color)
    add_text_box(slide1, tag, 4.25 + i * 1.82, 3.68, 1.6, 0.33,
                 font_size=11, bold=True, color=WHITE if tag != "JavaScript" else NAVY,
                 align=PP_ALIGN.CENTER)

# Subtítulo inferior
add_text_box(slide1, "Proyecto Final — Curso JavaScript Frontend", 4.2, 4.25, 9, 0.5,
             font_size=16, bold=False, color=RGBColor(0x88, 0x88, 0xaa),
             align=PP_ALIGN.LEFT)

# Nombre alumno / fecha
add_rect(slide1, 0.18, 6.8, 13.15, 0.7, BLUE)
add_text_box(slide1, "Mattias Maciel  |  2026", 0.5, 6.86, 12, 0.5,
             font_size=15, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# =============================================
#  SLIDE 2 — PROYECTO / CONTEXTO
# =============================================
slide2 = prs.slides.add_slide(BLANK)
add_rect(slide2, 0, 0, 13.33, 7.5, LIGHT_BG)
slide_dark_header(slide2, "Proyecto / Contexto",
                  "¿Qué es TechStore y qué problema resuelve?")

# Bloque descripción
add_rect(slide2, 0.45, 1.75, 5.8, 4.7, WHITE)
add_rect(slide2, 0.45, 1.75, 5.8, 0.06, RED)
add_text_box(slide2, "¿Qué es?", 0.65, 1.82, 5.4, 0.5,
             font_size=16, bold=True, color=NAVY)
add_text_box(slide2,
    "TechStore es un simulador interactivo de tienda "
    "de tecnología que replica el flujo completo de "
    "una compra online, desde el catálogo hasta la "
    "confirmación del pedido, utilizando únicamente "
    "JavaScript vanilla, HTML y CSS.",
    0.65, 2.3, 5.4, 2.0,
    font_size=13.5, color=GRAY_TEXT)

add_text_box(slide2, "Proceso de negocio simulado", 0.65, 3.85, 5.4, 0.5,
             font_size=14, bold=True, color=NAVY)
add_bullet_box(slide2,
    ["Catálogo con filtros por categoría",
     "Carrito lateral + página de carrito",
     "Formulario de envío con validación",
     "Selección de método de pago",
     "Resumen y confirmación del pedido",
     "Pantalla de compra exitosa"],
    0.65, 4.28, 5.4, 2.0,
    font_size=12.5, color=GRAY_TEXT)

# Estructura del proyecto
add_rect(slide2, 6.55, 1.75, 6.3, 2.65, WHITE)
add_rect(slide2, 6.55, 1.75, 6.3, 0.06, BLUE)
add_text_box(slide2, "Estructura del proyecto", 6.75, 1.82, 5.9, 0.5,
             font_size=15, bold=True, color=NAVY)
add_code_block(slide2,
    ["ecommerce-js/",
     "├── index.html",
     "├── css/",
     "│   └── style.css",
     "├── js/",
     "│   └── app.js",
     "└── data/",
     "    └── products.json"],
    6.55, 2.35, 6.3, 2.0, font_size=12)

# Flujo de pantallas (numerado)
add_rect(slide2, 6.55, 4.5, 6.3, 1.9, WHITE)
add_rect(slide2, 6.55, 4.5, 6.3, 0.06, RED)
add_text_box(slide2, "Flujo de pantallas (SPA)", 6.75, 4.57, 5.9, 0.45,
             font_size=15, bold=True, color=NAVY)

steps = ["1 Productos", "2 Carrito", "3 Datos", "4 Pago", "5 Confirmar", "6 Éxito"]
colors_flow = [BLUE, BLUE, BLUE, BLUE, BLUE, RGBColor(0x4c, 0xaf, 0x50)]
for i, (step, col) in enumerate(zip(steps, colors_flow)):
    col_x = 6.75 + i * 0.98
    add_rect(slide2, col_x, 5.1, 0.88, 0.38, col)
    add_text_box(slide2, step, col_x + 0.01, 5.12, 0.87, 0.34,
                 font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text_box(slide2, "→", col_x + 0.87, 5.13, 0.12, 0.34,
                     font_size=10, bold=True, color=GRAY_TEXT)

# =============================================
#  SLIDE 3 — CÓDIGO HTML
# =============================================
slide3 = prs.slides.add_slide(BLANK)
add_rect(slide3, 0, 0, 13.33, 7.5, LIGHT_BG)
slide_dark_header(slide3, "Código — HTML", "index.html  ·  SPA de 6 secciones + panel lateral")

# Columna izquierda: estructura HTML en código
add_text_box(slide3, "Estructura principal", 0.45, 1.72, 7.0, 0.45,
             font_size=14, bold=True, color=NAVY)
add_code_block(slide3,
    ['<!DOCTYPE html>',
     '<html lang="es">',
     '<head>',
     '  <!-- Toastify CSS (librería externa) -->',
     '  <link rel="stylesheet" href=".../toastify.min.css">',
     '  <link rel="stylesheet" href="css/style.css">',
     '</head>',
     '<body>',
     '  <header class="header">',
     '    <div id="barra-pasos" class="barra-pasos">',
     '      <!-- Pasos 1→4 del checkout -->',
     '    </div>',
     '  </header>',
     '  <main class="main">',
     '    <section id="pagina-productos" class="pagina activa">',
     '    <section id="pagina-carrito"   class="pagina">',
     '    <section id="pagina-datos"     class="pagina">',
     '    <section id="pagina-pago"      class="pagina">',
     '    <section id="pagina-confirmacion" class="pagina">',
     '    <section id="pagina-exito"     class="pagina">',
     '  </main>',
     '  <aside id="panel-carrito" class="panel-carrito">',
     '  <!-- Toastify JS -->',
     '  <script src=".../toastify.js"></script>',
     '  <script src="js/app.js"></script>',
     '</body>'],
    0.45, 2.15, 7.3, 5.0, font_size=10.5)

# Columna derecha: puntos clave
add_rect(slide3, 8.1, 1.72, 4.75, 5.43, WHITE)
add_rect(slide3, 8.1, 1.72, 4.75, 0.06, RED)
add_text_box(slide3, "Puntos clave", 8.3, 1.79, 4.4, 0.45,
             font_size=15, bold=True, color=NAVY)
add_bullet_box(slide3,
    ["SPA: todas las secciones en un único HTML",
     "Navegación por clases CSS (sin recarga)",
     "Barra de pasos visible solo en checkout",
     "Panel lateral del carrito como <aside>",
     "Toastify cargado desde CDN (reemplaza alert/confirm/prompt)",
     "JS al final del body para asegurar DOM cargado"],
    8.3, 2.32, 4.4, 3.5,
    font_size=12.5, color=GRAY_TEXT)

# Etiqueta técnica
add_rect(slide3, 8.3, 6.2, 2.2, 0.38, NAVY)
add_text_box(slide3, "6 páginas virtuales", 8.37, 6.23, 2.1, 0.33,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide3, 10.65, 6.2, 2.0, 0.38, BLUE)
add_text_box(slide3, "1 archivo HTML", 10.72, 6.23, 1.88, 0.33,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# =============================================
#  SLIDE 4 — CÓDIGO CSS
# =============================================
slide4 = prs.slides.add_slide(BLANK)
add_rect(slide4, 0, 0, 13.33, 7.5, LIGHT_BG)
slide_dark_header(slide4, "Código — CSS", "style.css  ·  Diseño responsivo y sistema de páginas")

# Sección 1: sistema de páginas
add_text_box(slide4, "Sistema de páginas virtuales", 0.45, 1.72, 6.0, 0.45,
             font_size=14, bold=True, color=NAVY)
add_code_block(slide4,
    ['/* Todas las páginas ocultas por defecto */',
     '.pagina {',
     '    display: none;',
     '}',
     '',
     '/* Solo la activa se muestra */',
     '.pagina.activa {',
     '    display: block;',
     '}',
     '',
     '/* Panel lateral con animación */',
     '.panel-carrito {',
     '    position: fixed;',
     '    right: -400px;          /* Oculto fuera de pantalla */',
     '    transition: right 0.3s ease;',
     '}',
     '.panel-carrito.abierto {',
     '    right: 0;               /* Desliza al interior */',
     '}'],
    0.45, 2.15, 6.0, 4.25, font_size=10.5)

# Sección 2: grilla responsive
add_text_box(slide4, "Grilla de productos responsive", 6.75, 1.72, 6.1, 0.45,
             font_size=14, bold=True, color=NAVY)
add_code_block(slide4,
    ['.productos-grid {',
     '    display: grid;',
     '    grid-template-columns:',
     '        repeat(auto-fill, minmax(260px, 1fr));',
     '    gap: 1.5rem;',
     '}',
     '',
     '@media (max-width: 700px) {',
     '    .productos-grid {',
     '        grid-template-columns: 1fr;',
     '    }',
     '    .form-fila {',
     '        grid-template-columns: 1fr;',
     '    }',
     '}'],
    6.75, 2.15, 6.1, 3.0, font_size=10.5)

# Paleta de colores
add_rect(slide4, 6.75, 5.3, 6.1, 1.85, WHITE)
add_rect(slide4, 6.75, 5.3, 6.1, 0.06, RED)
add_text_box(slide4, "Paleta de colores", 6.95, 5.37, 5.7, 0.45,
             font_size=14, bold=True, color=NAVY)
palette = [
    ("#1a1a2e", "Fondo/Header", NAVY),
    ("#0f3460", "Botones/Primary", BLUE),
    ("#e94560", "Acento/Precio", RED),
    ("#f5f5f5", "Fondo contenido", LIGHT_BG),
]
for i, (hex_val, label, col) in enumerate(palette):
    bx = 6.95 + i * 1.47
    add_rect(slide4, bx, 5.9, 1.3, 0.5, col)
    add_text_box(slide4, hex_val, bx, 6.45, 1.3, 0.3,
                 font_size=9, bold=True, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    add_text_box(slide4, label, bx, 6.75, 1.3, 0.3,
                 font_size=9, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# =============================================
#  SLIDE 5 — JS: FETCH Y ARRAYS
# =============================================
slide5 = prs.slides.add_slide(BLANK)
add_rect(slide5, 0, 0, 13.33, 7.5, LIGHT_BG)
slide_dark_header(slide5, "Código — JavaScript: Fetch y Arrays",
                  "app.js  ·  Carga de datos desde JSON externo + métodos funcionales")

# Código Fetch
add_text_box(slide5, "Carga de productos con Fetch API", 0.45, 1.72, 7.2, 0.45,
             font_size=14, bold=True, color=NAVY)
add_code_block(slide5,
    ['function cargarProductos() {',
     '    fetch("data/products.json")',
     '        .then(function (respuesta) {',
     '            if (!respuesta.ok) {',
     '                throw new Error("Error al cargar productos");',
     '            }',
     '            return respuesta.json();',
     '        })',
     '        .then(function (datos) {',
     '            productos = datos;',
     '',
     '            // Ordenar por categoría y luego por precio',
     '            productos.sort(function (a, b) {',
     '                if (a.categoria < b.categoria) return -1;',
     '                if (a.categoria > b.categoria) return 1;',
     '                return a.precio - b.precio;',
     '            });',
     '',
     '            mostrarProductos(productos);',
     '        })',
     '        .catch(function () {',
     '            cargandoDiv.innerHTML = "Error al cargar productos.";',
     '        });',
     '}'],
    0.45, 2.15, 7.2, 4.8, font_size=10.5)

# Métodos de arrays
add_rect(slide5, 7.95, 1.72, 4.9, 5.43, WHITE)
add_rect(slide5, 7.95, 1.72, 4.9, 0.06, RED)
add_text_box(slide5, "Métodos de arrays utilizados", 8.15, 1.79, 4.55, 0.45,
             font_size=14, bold=True, color=NAVY)

array_methods = [
    (".sort()",   "Ordena por categoría y precio al cargar"),
    (".filter()", "Filtra productos por categoría"),
    (".find()",   "Busca producto por ID en el carrito"),
    (".reduce()","Calcula el total de la compra"),
    (".forEach()","Itera productos y carrito para renderizar"),
]
for i, (method, desc) in enumerate(array_methods):
    ypos = 2.38 + i * 0.9
    add_rect(slide5, 8.15, ypos, 1.4, 0.4, NAVY)
    add_text_box(slide5, method, 8.18, ypos + 0.04, 1.35, 0.34,
                 font_size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text_box(slide5, desc, 9.65, ypos + 0.05, 2.95, 0.35,
                 font_size=11.5, color=GRAY_TEXT)

# Miniatura del JSON
add_rect(slide5, 7.95, 6.55, 4.9, 0.6, CODE_BG)
add_text_box(slide5,
    '// data/products.json  →  array de 13 productos',
    8.1, 6.6, 4.6, 0.45,
    font_size=10.5, color=RGBColor(0x88, 0xcc, 0xff),
    italic=True)

# =============================================
#  SLIDE 6 — JS: DOM Y EVENTOS
# =============================================
slide6 = prs.slides.add_slide(BLANK)
add_rect(slide6, 0, 0, 13.33, 7.5, LIGHT_BG)
slide_dark_header(slide6, "Código — JavaScript: DOM y Eventos",
                  "app.js  ·  Navegación entre páginas y validación de formularios")

# Función navegarA
add_text_box(slide6, "Función central: navegarA()", 0.45, 1.72, 7.2, 0.45,
             font_size=14, bold=True, color=NAVY)
add_code_block(slide6,
    ['function navegarA(nombrePagina) {',
     '',
     '    // Ocultar todas las páginas',
     '    paginas.forEach(function (pagina) {',
     '        pagina.classList.remove("activa");',
     '    });',
     '',
     '    // Mostrar la página destino',
     '    var paginaDestino = document.getElementById(nombrePagina);',
     '    if (paginaDestino) {',
     '        paginaDestino.classList.add("activa");',
     '    }',
     '',
     '    // Mostrar/ocultar barra de pasos',
     '    var paginasCheckout = [',
     '        "pagina-carrito","pagina-datos",',
     '        "pagina-pago","pagina-confirmacion"',
     '    ];',
     '    if (paginasCheckout.indexOf(nombrePagina) !== -1) {',
     '        barraPasos.classList.add("visible");',
     '    } else {',
     '        barraPasos.classList.remove("visible");',
     '    }',
     '    actualizarBarraPasos(nombrePagina);',
     '    cerrarPanelCarrito();',
     '}'],
    0.45, 2.15, 7.2, 4.8, font_size=10.5)

# Validación de formulario
add_rect(slide6, 7.95, 1.72, 4.9, 5.43, WHITE)
add_rect(slide6, 7.95, 1.72, 4.9, 0.06, RED)
add_text_box(slide6, "Validación en tiempo real", 8.15, 1.79, 4.55, 0.45,
             font_size=14, bold=True, color=NAVY)
add_code_block(slide6,
    ['// Validación en vivo del email',
     'emailInput.addEventListener("blur",',
     '    function () {',
     '        validarEmailEnVivo();',
     '    }',
     ');',
     '',
     'emailInput.addEventListener("input",',
     '    function () {',
     '        if (emailInput.classList',
     '               .contains("error-input")) {',
     '            validarEmailEnVivo();',
     '        }',
     '    }',
     ');'],
    7.95, 2.28, 4.9, 2.8, font_size=10.5)

add_text_box(slide6, "Eventos utilizados", 8.15, 5.22, 4.55, 0.4,
             font_size=14, bold=True, color=NAVY)
add_bullet_box(slide6,
    ["click  — botones de navegación y carrito",
     "submit — formularios de datos y pago",
     "blur / input — validación en tiempo real",
     "change — cambio de método de pago"],
    8.15, 5.65, 4.55, 1.5,
    font_size=11.5, color=GRAY_TEXT)

# =============================================
#  SLIDE 7 — LIBRERÍA TOASTIFY
# =============================================
slide7 = prs.slides.add_slide(BLANK)
add_rect(slide7, 0, 0, 13.33, 7.5, LIGHT_BG)
slide_dark_header(slide7, "Librería Externa — Toastify JS",
                  "Notificaciones toast en reemplazo de alert(), confirm() y prompt()")

# Columna izquierda: código
add_text_box(slide7, "Implementación en el proyecto", 0.45, 1.72, 6.5, 0.45,
             font_size=14, bold=True, color=NAVY)
add_code_block(slide7,
    ['<!-- Carga desde CDN en index.html -->',
     '<script src="https://cdn.jsdelivr.net/',
     '   npm/toastify-js"></script>',
     '',
     '// Función wrapper en app.js',
     'function mostrarToast(mensaje, tipo) {',
     '    var colores = {',
     '        ok:    "#4caf50",   // verde',
     '        error: "#e94560",   // rojo',
     '        info:  "#0f3460"    // azul',
     '    };',
     '',
     '    Toastify({',
     '        text:     mensaje,',
     '        duration: 2500,',
     '        gravity:  "bottom",',
     '        position: "right",',
     '        style: {',
     '            background: colores[tipo] || colores.info',
     '        }',
     '    }).showToast();',
     '}'],
    0.45, 2.15, 6.5, 4.8, font_size=10.5)

# Columna derecha
add_rect(slide7, 7.2, 1.72, 5.65, 2.5, WHITE)
add_rect(slide7, 7.2, 1.72, 5.65, 0.06, RED)
add_text_box(slide7, "¿Por qué Toastify?", 7.4, 1.79, 5.2, 0.45,
             font_size=15, bold=True, color=NAVY)
add_bullet_box(slide7,
    ["Elimina alert(), confirm() y prompt() (requisito del curso)",
     "Notificaciones no bloqueantes → mejor UX",
     "3 tipos: ok (verde), error (rojo), info (azul)",
     "Posición y duración configurables",
     "Integración con una sola función wrapper"],
    7.4, 2.28, 5.2, 1.9,
    font_size=12, color=GRAY_TEXT)

# Usos en el proyecto
add_rect(slide7, 7.2, 4.35, 5.65, 2.8, WHITE)
add_rect(slide7, 7.2, 4.35, 5.65, 0.06, BLUE)
add_text_box(slide7, "Dónde se usa mostrarToast()", 7.4, 4.42, 5.2, 0.45,
             font_size=14, bold=True, color=NAVY)
toast_uses = [
    ("ok",    "Producto agregado al carrito"),
    ("ok",    "Pedido confirmado con éxito"),
    ("error", "Sin stock del producto"),
    ("error", "Carrito vacío"),
    ("error", "Campos del formulario incompletos"),
    ("info",  "Carrito vaciado"),
    ("info",  "Producto eliminado del carrito"),
]
cols_toast = {"ok": RGBColor(0x4c, 0xaf, 0x50), "error": RED, "info": BLUE}
for i, (tipo, msg) in enumerate(toast_uses):
    yp = 4.95 + i * 0.28
    add_rect(slide7, 7.4, yp, 0.55, 0.22, cols_toast[tipo])
    add_text_box(slide7, tipo, 7.41, yp + 0.01, 0.54, 0.2,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide7, msg, 8.02, yp, 4.5, 0.25,
                 font_size=11, color=GRAY_TEXT)

# =============================================
#  SLIDE 8 — REPOSITORIO
# =============================================
slide8 = prs.slides.add_slide(BLANK)
add_rect(slide8, 0, 0, 13.33, 7.5, NAVY)
add_rect(slide8, 0, 0, 0.18, 7.5, RED)
add_rect(slide8, 0.18, 0, 13.15, 0.1, BLUE)

add_text_box(slide8, "Repositorio", 1.0, 0.6, 11, 1.0,
             font_size=42, bold=True, color=WHITE)
add_rect(slide8, 1.0, 1.65, 6.0, 0.08, RED)

# Ícono GitHub
add_text_box(slide8, "⌥", 1.0, 2.1, 1.5, 1.5,
             font_size=72, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide8, 2.8, 2.2, 9.5, 0.72, BLUE)
add_text_box(slide8, "🔗  github.com / [tu-usuario] / ecommerce-js",
             3.0, 2.28, 9.1, 0.55,
             font_size=19, bold=True, color=WHITE)

add_text_box(slide8,
    "⚠  Completar con el link real del repositorio antes de entregar",
    2.8, 3.05, 9.5, 0.5,
    font_size=13, italic=True, color=RGBColor(0xff, 0xdd, 0x57))

# Datos del repo
add_rect(slide8, 1.0, 3.75, 10.8, 2.8, BLUE)
for i, (label, value) in enumerate([
    ("Rama principal", "main"),
    ("Lenguajes",      "HTML  ·  CSS  ·  JavaScript"),
    ("Archivos clave", "index.html  ·  js/app.js  ·  css/style.css  ·  data/products.json"),
    ("Dependencias",   "Ninguna (solo CDN de Toastify JS)"),
]):
    yp = 3.95 + i * 0.6
    add_text_box(slide8, label + ":", 1.2, yp, 2.8, 0.45,
                 font_size=13, bold=True, color=RGBColor(0xcc, 0xdd, 0xff))
    add_text_box(slide8, value, 4.1, yp, 7.4, 0.45,
                 font_size=13, color=WHITE)

# =============================================
#  SLIDE 9 — DEMO
# =============================================
slide9 = prs.slides.add_slide(BLANK)
add_rect(slide9, 0, 0, 13.33, 7.5, LIGHT_BG)
slide_dark_header(slide9, "Demo — Flujo de la aplicación",
                  "Capturas del simulador en funcionamiento")

demo_screens = [
    ("🏪", "Catálogo",     "Productos cargados\ndesde JSON vía Fetch.\nFiltros por categoría."),
    ("🛒", "Carrito",      "Panel lateral\ndeslizante con\nanimación CSS."),
    ("📋", "Datos",        "Formulario con\nvalidación en\ntiempo real."),
    ("💳", "Pago",         "4 métodos de pago.\nCampos de tarjeta\ndinámicos."),
    ("✅", "Confirmación", "Resumen completo\ndel pedido antes\nde confirmar."),
    ("🎉", "Éxito",        "Mensaje final con\nnombre y dirección\ndel cliente."),
]

for i, (emoji, title, desc) in enumerate(demo_screens):
    col = i % 3
    row = i // 3
    bx = 0.45 + col * 4.27
    by = 1.75 + row * 2.75

    add_rect(slide9, bx, by, 4.0, 2.4, WHITE)
    add_rect(slide9, bx, by, 4.0, 0.06, RED if row == 0 else BLUE)
    add_text_box(slide9, emoji, bx + 0.1, by + 0.1, 0.8, 0.8,
                 font_size=32, align=PP_ALIGN.CENTER)
    add_text_box(slide9, title, bx + 0.95, by + 0.1, 2.9, 0.45,
                 font_size=16, bold=True, color=NAVY)
    add_text_box(slide9, desc, bx + 0.15, by + 0.65, 3.7, 1.5,
                 font_size=12, color=GRAY_TEXT)

add_text_box(slide9,
    "💡  Agregar capturas de pantalla reales sobre cada tarjeta antes de entregar",
    0.45, 7.1, 12.5, 0.4,
    font_size=11, italic=True,
    color=RGBColor(0x88, 0x88, 0xaa))

# =============================================
#  SLIDE 10 — REFLEXIÓN FINAL (placeholder)
# =============================================
slide10 = prs.slides.add_slide(BLANK)
add_rect(slide10, 0, 0, 13.33, 7.5, NAVY)
add_rect(slide10, 0, 0, 0.18, 7.5, RED)
add_rect(slide10, 0.18, 0, 13.15, 0.1, BLUE)

add_text_box(slide10, "Reflexión Final", 1.0, 0.6, 11, 1.0,
             font_size=42, bold=True, color=WHITE)
add_rect(slide10, 1.0, 1.65, 6.0, 0.08, RED)

add_rect(slide10, 1.0, 2.1, 11, 4.5, BLUE)
add_text_box(slide10, "✏  Esta sección se completa próximamente", 1.3, 3.4, 10.5, 0.9,
             font_size=22, bold=True, color=RGBColor(0xff, 0xdd, 0x57),
             align=PP_ALIGN.CENTER)
add_text_box(slide10,
    "Pendiente de redacción.\nAquí irá la reflexión sobre el proceso de aprendizaje,\n"
    "los desafíos encontrados y los conocimientos adquiridos.",
    1.3, 4.1, 10.5, 2.0,
    font_size=16, color=RGBColor(0xcc, 0xdd, 0xff),
    align=PP_ALIGN.CENTER)

# =============================================
#  GUARDAR
# =============================================
output_path = r"C:\Users\matti\Documents\ecommerce-js\TechStore_Presentacion.pptx"
prs.save(output_path)
print(f"Archivo guardado: {output_path}")
